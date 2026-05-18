"""build_outputs.py — PDF → PPTX（每页截图铺满）"""
import sys
from pathlib import Path
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches, Emu
import tempfile

pdf_path  = sys.argv[1]
pptx_path = sys.argv[2]

# PDF → PNG（150 DPI，足够清晰）
print(f"  转换 PDF 页面…")
pages = convert_from_path(pdf_path, dpi=150)
print(f"  共 {len(pages)} 页")

# PPTX：每页铺一张全尺寸图
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

with tempfile.TemporaryDirectory() as tmpdir:
    for i, page in enumerate(pages):
        img_path = f"{tmpdir}/p{i:03d}.png"
        page.save(img_path, "PNG")
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(img_path, Emu(0), Emu(0), prs.slide_width, prs.slide_height)
        sys.stdout.write(f"  幻灯片 {i+1:02d}/{len(pages)}\r")

print()
prs.save(pptx_path)
print(f"  ✅ PPTX: {len(pages)} 页 → {pptx_path}")
