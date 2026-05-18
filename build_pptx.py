#!/usr/bin/env python3
"""
build_pptx.py — 从 Marp MD 生成可编辑 PPTX 草稿

结构：
  pptx-export/
  ├── dual-mode/
  │   ├── background.png      ← 1 张深色渐变背景（无文字）
  │   ├── content/            ← slide_01.md … slide_27.md（每页文字内容）
  │   ├── images/             ← 图片资产（logo、截图等）
  │   └── slides_dual_v3.pptx
  └── chatflow/
      ├── background.png
      ├── content/
      ├── images/
      └── slides_chatflow_v3.pptx

用法:
  python3 build_pptx.py dual        # 双模版
  python3 build_pptx.py chatflow    # Chatflow版
  python3 build_pptx.py all         # 两个都生成
  python3 build_pptx.py all --force # 强制重新提取（覆盖已有文件）
"""

import re
import sys
from pathlib import Path

from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu, Inches, Pt

BASE = Path(__file__).parent

CONFIGS = {
    "dual": {
        "md":    BASE / "路演PPT-v2.md",
        "outdir": BASE / "pptx-export" / "dual-mode",
        "pptx":   BASE / "pptx-export" / "dual-mode" / "slides_dual_v3.pptx",
        "label":  "双模版",
    },
    "chatflow": {
        "md":    BASE / "路演PPT-v2-chatflow.md",
        "outdir": BASE / "pptx-export" / "chatflow",
        "pptx":   BASE / "pptx-export" / "chatflow" / "slides_chatflow_v3.pptx",
        "label":  "Chatflow版",
    },
}

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# 主题色 #0b1527 → #0d1b2a（竖向渐变）
BG_TOP    = (11, 21, 39)
BG_BOTTOM = (13, 27, 42)

# 文字颜色
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_H2     = RGBColor(0xE2, 0xE8, 0xF0)
C_BODY   = RGBColor(0xCB, 0xD5, 0xE1)
C_SUBTLE = RGBColor(0x8A, 0x9B, 0xB8)


# ── 背景图生成 ──────────────────────────────────────────────────

def create_background(path: Path, width: int = 1920, height: int = 1080) -> Path:
    """生成一张纯色渐变深色背景（无文字）"""
    img = Image.new("RGB", (width, height))
    draw = ImageDraw.Draw(img)
    for y in range(height):
        t = y / height
        r = int(BG_TOP[0] + (BG_BOTTOM[0] - BG_TOP[0]) * t)
        g = int(BG_TOP[1] + (BG_BOTTOM[1] - BG_TOP[1]) * t)
        b = int(BG_TOP[2] + (BG_BOTTOM[2] - BG_TOP[2]) * t)
        draw.line([(0, y), (width, y)], fill=(r, g, b))
    path.parent.mkdir(parents=True, exist_ok=True)
    img.save(str(path))
    print(f"  ✅ 背景图已生成 → {path}")
    return path


# ── MD 解析 ─────────────────────────────────────────────────────

def _strip_html(text: str) -> str:
    text = re.sub(r'<[^>]+>', ' ', text)
    text = re.sub(r'\s+', ' ', text)
    return text.strip()


def _parse_slides(md_path: Path) -> list[str]:
    """把 Marp MD 按 --- 切成 slide 块，返回每块的纯文字 MD 字符串。"""
    raw = md_path.read_text(encoding="utf-8")
    # 去掉 YAML frontmatter
    raw = re.sub(r'^---.*?---\n', '', raw, count=1, flags=re.DOTALL)
    blocks = re.split(r'\n---\n', raw)
    results = []
    for block in blocks:
        results.append(_extract_readable(block))
    return results


def _extract_readable(block: str) -> str:
    """从含 HTML 的 slide 块提取可读文字，保留 Markdown 结构。"""
    # 先整块移除 <style>...</style> 和 base64 data URI（避免 CSS 污染）
    block = re.sub(r'<style[^>]*>.*?</style>', '', block, flags=re.DOTALL | re.IGNORECASE)
    block = re.sub(r'data:[^"\')\s]+', '[img]', block)
    lines_out = []
    for line in block.splitlines():
        s = line.strip()
        if not s:
            continue
        # 跳过注释和 Marp 指令
        if s.startswith('<!--') or s.startswith('-->') or s.startswith('<!-- _'):
            continue
        # 纯 Markdown 标题
        m = re.match(r'^(#{1,4}) (.+)', s)
        if m:
            text = _strip_html(m.group(2))
            if text:
                lines_out.append(f"{m.group(1)} {text}")
            continue
        # 列表项
        m = re.match(r'^[-*] (.+)', s)
        if m:
            text = _strip_html(m.group(1))
            if text:
                lines_out.append(f"- {text}")
            continue
        # 引用
        if s.startswith('> '):
            text = _strip_html(s[2:])
            if text:
                lines_out.append(f"> {text}")
            continue
        # 表格行（跳过分隔行）
        if s.startswith('|') and not re.match(r'^\|[-:| ]+\|$', s):
            cols = [c.strip() for c in s.split('|') if c.strip()]
            if cols:
                lines_out.append('  '.join(cols))
            continue
        # HTML 行 → 提取文字
        if s.startswith('<'):
            text = _strip_html(s)
            if text and len(text) > 3 and not re.match(r'^[\d.%→✅❌⚠️]+$', text):
                lines_out.append(text)
            continue
        # 普通文字（排除 style= 等属性碎片）
        if len(s) > 3 and not s.startswith('style=') and not s.startswith('class='):
            lines_out.append(s)

    return '\n'.join(lines_out)


def save_slide_content(slides: list[str], content_dir: Path) -> None:
    """把每张 slide 的文字内容写入 content/slide_XX.md"""
    content_dir.mkdir(parents=True, exist_ok=True)
    for i, content in enumerate(slides, 1):
        path = content_dir / f"slide_{i:02d}.md"
        path.write_text(content, encoding="utf-8")
    print(f"  ✅ {len(slides)} 张幻灯片内容 → {content_dir}")


# ── PPTX 组装 ───────────────────────────────────────────────────

def _add_text_run(para, text: str, size: int, bold: bool, color: RGBColor) -> None:
    run = para.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def build_pptx(bg_path: Path, slides: list[str], pptx_path: Path) -> None:
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    for i, content in enumerate(slides, 1):
        slide = prs.slides.add_slide(blank)

        # 背景图（全页铺满）
        slide.shapes.add_picture(
            str(bg_path),
            Emu(0), Emu(0),
            prs.slide_width,
            prs.slide_height,
        )

        # 文字框（内容区域）
        txBox = slide.shapes.add_textbox(
            Inches(0.8), Inches(0.4),
            Inches(11.7), Inches(6.8),
        )
        tf = txBox.text_frame
        tf.word_wrap = True

        first_para = True
        for line in content.splitlines():
            s = line.strip()
            if not s:
                continue
            p = tf.paragraphs[0] if first_para else tf.add_paragraph()
            first_para = False

            if re.match(r'^# ', s):
                _add_text_run(p, s[2:], 40, True, C_WHITE)
                p.space_after = Pt(10)
            elif re.match(r'^## ', s):
                _add_text_run(p, s[3:], 28, True, C_H2)
                p.space_after = Pt(8)
            elif re.match(r'^### ', s):
                _add_text_run(p, s[4:], 22, True, C_H2)
            elif s.startswith('- '):
                _add_text_run(p, f"• {s[2:]}", 18, False, C_BODY)
            elif s.startswith('> '):
                _add_text_run(p, s[2:], 16, False, C_SUBTLE)
            else:
                _add_text_run(p, s, 16, False, C_BODY)

        print(f"  幻灯片 {i:02d}/{len(slides)} ✓", end="\r")

    print()
    pptx_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(pptx_path))
    print(f"  ✅ PPTX 已保存 → {pptx_path}")


# ── 主流程 ──────────────────────────────────────────────────────

def run(key: str, force: bool = False) -> None:
    cfg = CONFIGS[key]
    print(f"\n{'='*52}")
    print(f"  生成{cfg['label']} PPTX")
    print(f"{'='*52}")

    if not cfg["md"].exists():
        print(f"  ❌ 找不到 MD 文件：{cfg['md']}")
        return

    out = cfg["outdir"]
    bg_path = out / "background.png"
    content_dir = out / "content"

    # 背景图（只需生成一次）
    if not bg_path.exists() or force:
        create_background(bg_path)
    else:
        print("  ⏩ 背景图已存在，跳过（加 --force 重新生成）")

    # 解析 slide 内容
    print("  📄 解析幻灯片文字内容…")
    slides = _parse_slides(cfg["md"])
    print(f"     共 {len(slides)} 张幻灯片")
    save_slide_content(slides, content_dir)

    # 组装 PPTX
    print("  📊 组装 PPTX…")
    build_pptx(bg_path, slides, cfg["pptx"])


if __name__ == "__main__":
    args = sys.argv[1:]
    force = "--force" in args
    keys = [a for a in args if not a.startswith("--")]

    if not keys or keys[0] not in ("dual", "chatflow", "all"):
        print(__doc__)
        sys.exit(1)

    targets = ["dual", "chatflow"] if keys[0] == "all" else [keys[0]]
    for k in targets:
        run(k, force=force)
